from pptx import Presentation
import glob

# code snippet to read content from *.pptx and save to a *.txt file. 
# Ref: https://stackoverflow.com/questions/13559133/how-to-open-ppt-file-using-python#:~:text=If%20you%20have%20access%20to,api%20to%20read%20the%20file.&text=python%2Dpptx%20can%20open%20recent,in%20their%20Getting%20started%20guide.&text=Using%20catdoc%2Fcatppt%20with%20subprocess,doc%20files%20and%20ppt%20files.
# Step 1: Install package python-pptx in python vern
# Step 2: Create a name.txt file where contents to save into
# Step 3: Copy *.pptx files to the same directory
# Step 4: run below and done!

words = open("words.txt", "w")
for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # print(shape.text)
                words.write(shape.text)